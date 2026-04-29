#!/usr/bin/env python3
"""
智能 PPT 生成器 - 根据主题自动判断设计风格和元素
输入：主题描述 + 内容
输出：匹配风格的 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

# ==================== 配色方案库 ====================
class ColorPalette:
    """配色方案"""
    
    # 疗愈温柔风
    HEALING = {
        "name": "疗愈温柔风",
        "primary": RGBColor(0x9B, 0x8A, 0xB8),      # 薰衣草紫
        "secondary": RGBColor(0xF7, 0xD4, 0xC4),    # 藕粉
        "accent": RGBColor(0xA8, 0xD8, 0xEA),       # 薄荷绿
        "light": RGBColor(0xF9, 0xF5, 0xF2),        # 米白
        "text": RGBColor(0x4A, 0x4A, 0x4A),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
    }
    
    # 商务专业风
    BUSINESS = {
        "name": "商务专业风",
        "primary": RGBColor(0x1A, 0x37, 0x62),      # 深藏蓝
        "secondary": RGBColor(0x4A, 0x90, 0xD9),   # 亮蓝
        "accent": RGBColor(0xF5, 0xA6, 0x23),      # 金色
        "light": RGBColor(0xE8, 0xF4, 0xF8),
        "text": RGBColor(0x2C, 0x3E, 0x50),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
    }
    
    # 活泼温暖风
    LIVELY = {
        "name": "活泼温暖风",
        "primary": RGBColor(0xFF, 0x6B, 0x6B),      # 珊瑚红
        "secondary": RGBColor(0xFF, 0x9F, 0x43),   # 活力橙
        "accent": RGBColor(0x00, 0xD2, 0xFF),      # 天蓝
        "light": RGBColor(0xFF, 0xF3, 0xF0),
        "text": RGBColor(0x2D, 0x34, 0x36),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
    }
    
    # 自然清新风
    NATURAL = {
        "name": "自然清新风",
        "primary": RGBColor(0x2D, 0x86, 0x6C),     # 森林绿
        "secondary": RGBColor(0xA8, 0xD8, 0xEA),   # 薄荷蓝
        "accent": RGBColor(0xF7, 0xD4, 0xC4),      # 沙色
        "light": RGBColor(0xF0, 0xF7, 0xF4),
        "text": RGBColor(0x2D, 0x34, 0x36),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
    }
    
    # 神秘深邃风
    MYSTIC = {
        "name": "神秘深邃风",
        "primary": RGBColor(0x2C, 0x3E, 0x50),      # 深灰蓝
        "secondary": RGBColor(0x9B, 0x59, 0xB6),   # 紫色
        "accent": RGBColor(0xF1, 0xC4, 0x0C),      # 金色
        "light": RGBColor(0x34, 0x49, 0x5E),
        "text": RGBColor(0xE8, 0xE8, 0xE8),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
    }
    
    @classmethod
    def detect(cls, theme_text: str) -> dict:
        """根据主题关键词自动匹配配色方案"""
        theme_lower = theme_text.lower()
        
        # 关键词 → 配色方案映射
        keywords_map = {
            "疗愈": cls.HEALING,
            "温柔": cls.HEALING,
            "心理": cls.HEALING,
            "情感": cls.HEALING,
            "女性": cls.HEALING,
            "塔罗": cls.MYSTIC,
            "占星": cls.MYSTIC,
            "神秘": cls.MYSTIC,
            "欧卡": cls.HEALING,
            "冥想": cls.HEALING,
            "沙龙": cls.LIVELY,
            "活泼": cls.LIVELY,
            "温暖": cls.LIVELY,
            "女性成长": cls.LIVELY,
            "商务": cls.BUSINESS,
            "正式": cls.BUSINESS,
            "专业": cls.BUSINESS,
            "财富": cls.BUSINESS,
            "理财": cls.BUSINESS,
            "自然": cls.NATURAL,
            "清新": cls.NATURAL,
        }
        
        # 遍历匹配
        for keyword, palette in keywords_map.items():
            if keyword in theme_lower:
                return palette
        
        # 默认返回疗愈风
        return cls.HEALING


# ==================== 内容分析器 ====================
class ContentAnalyzer:
    """内容分析器 - 判断内容类型"""
    
    @staticmethod
    def analyze(content: dict) -> dict:
        """分析内容结构，返回推荐的设计元素"""
        result = {
            "type": "default",
            "elements": [],
            "layout": "cards",
            "icon_style": "circle",
        }
        
        content_text = str(content).lower()
        
        # 数据分析类
        if any(k in content_text for k in ["数据", "统计", "比例", "百分比", "占比", "chart", "数据"]):
            result["type"] = "data"
            result["elements"].append("chart")
            result["layout"] = "chart_focus"
        
        # 流程步骤类
        if any(k in content_text for k in ["步骤", "流程", "方法", "如何", "step", "流程"]):
            result["type"] = "process"
            result["elements"].append("flowchart")
            result["layout"] = "steps"
        
        # 人物介绍类
        if any(k in content_text for k in ["讲师", "导师", "介绍", "about", "bio"]):
            result["type"] = "person"
            result["elements"].append("photo_frame")
            result["layout"] = "person_focus"
        
        # 对比类
        if any(k in content_text for k in ["对比", "比较", "区别", "vs", "versus"]):
            result["type"] = "comparison"
            result["elements"].append("comparison_table")
            result["layout"] = "side_by_side"
        
        # 列表展示类
        if any(k in content_text for k in ["要点", "清单", "列表", "特点", "features"]):
            result["type"] = "list"
            result["elements"].append("icon_cards")
            result["layout"] = "grid"
        
        # 时间线类
        if any(k in content_text for k in ["时间", "历程", "发展", "timeline", "历史"]):
            result["type"] = "timeline"
            result["elements"].append("timeline")
            result["layout"] = "horizontal"
        
        return result


# ==================== 设计规则引擎 ====================
class DesignRules:
    """设计规则引擎"""
    
    @staticmethod
    def get_shape_for_type(element_type: str, palette: dict):
        """根据元素类型返回推荐形状"""
        shapes_map = {
            "step": MSO_SHAPE.ROUNDED_RECTANGLE,
            "card": MSO_SHAPE.ROUNDED_RECTANGLE,
            "highlight": MSO_SHAPE.OVAL,
            "flow": MSO_SHAPE.CHEVRON,
            "badge": MSO_SHAPE.PENTAGON,
        }
        return shapes_map.get(element_type, MSO_SHAPE.ROUNDED_RECTANGLE)
    
    @staticmethod
    def get_font_size_for_context(context: str) -> Pt:
        """根据上下文返回推荐字号"""
        sizes = {
            "title": Pt(36),
            "subtitle": Pt(24),
            "heading": Pt(20),
            "body": Pt(14),
            "caption": Pt(11),
            "footer": Pt(10),
        }
        return sizes.get(context, Pt(14))


# ==================== 智能 PPT 生成器 ====================
class SmartPPTGerator:
    """智能 PPT 生成器"""
    
    def __init__(self, theme: str, palette: dict = None):
        self.theme = theme
        self.palette = palette or ColorPalette.detect(theme)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def add_cover(self, title: str, subtitle: str = ""):
        """添加封面"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 渐变背景
        self._add_gradient_bg(slide, 0, 0, 13.333, 7.5, 
                              self.palette["primary"], self.palette["secondary"], 135)
        
        # 装饰圆形
        self._add_circle(slide, Inches(10.5), Inches(0.5), Inches(2), 
                        self.palette["accent"])
        self._add_circle(slide, Inches(11.5), Inches(5), Inches(1.5), 
                        self.palette["secondary"])
        
        # 主标题
        self._add_text(slide, title, Inches(1), Inches(2.5), 
                      Inches(11), Inches(1.2),
                      font_size=Pt(44), color=self.palette["white"], bold=True,
                      alignment=PP_ALIGN.CENTER)
        
        # 副标题
        if subtitle:
            self._add_text(slide, subtitle, Inches(1), Inches(4.2), 
                          Inches(11), Inches(0.8),
                          font_size=Pt(20), color=self.palette["light"],
                          alignment=PP_ALIGN.CENTER)
        
        # 风格标签
        self._add_text(slide, f"「{self.palette['name']}」", 
                      Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                      font_size=Pt(12), color=self.palette["light"],
                      alignment=PP_ALIGN.CENTER)
    
    def add_content_slide(self, title: str, content: dict):
        """添加内容页 - 自动判断元素类型"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 顶部装饰条
        self._add_gradient_bg(slide, 0, 0, 13.333, 0.8,
                              self.palette["primary"], self.palette["secondary"])
        
        # 标题
        self._add_text(slide, title, Inches(0.5), Inches(0.15), 
                      Inches(12), Inches(0.6),
                      font_size=Pt(28), color=self.palette["white"], bold=True)
        
        # 分析内容类型
        analysis = ContentAnalyzer.analyze(content)
        
        # 根据类型渲染内容
        if analysis["type"] == "data":
            self._render_data_content(slide, content)
        elif analysis["type"] == "process":
            self._render_process_content(slide, content)
        elif analysis["type"] == "list":
            self._render_list_content(slide, content)
        else:
            self._render_default_content(slide, content)
    
    def _render_data_content(self, slide, content: dict):
        """渲染数据类内容"""
        if "chart_data" in content:
            data = content["chart_data"]
            chart_data = CategoryChartData()
            chart_data.categories = data["categories"]
            for series_name, values in data["series"].items():
                chart_data.add_series(series_name, values)
            
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(2), Inches(1.5), Inches(9), Inches(5),
                chart_data
            )
            
            # 设置颜色
            chart = chart_shape.chart
            colors = [self.palette["primary"], self.palette["secondary"], 
                     self.palette["accent"]]
            for idx, series in enumerate(chart.series):
                if idx < len(colors):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = colors[idx]
    
    def _render_process_content(self, slide, content: dict):
        """渲染流程类内容"""
        steps = content.get("steps", [])
        if not steps:
            steps = list(content.values())[:5] if isinstance(content, dict) else []
        
        step_width = Inches(2.2)
        step_height = Inches(2.5)
        start_x = Inches(0.8)
        start_y = Inches(2.5)
        
        for idx, step in enumerate(steps):
            x = start_x + idx * (step_width + Inches(0.3))
            
            # 步骤卡片
            self._add_card(slide, x, start_y, step_width, step_height, 
                          self.palette["white"])
            
            # 步骤编号
            self._add_circle(slide, x + Inches(0.85), start_y + Inches(0.2), 
                           Inches(0.5), self.palette["primary"])
            self._add_text(slide, str(idx + 1), 
                         x + Inches(0.85), start_y + Inches(0.25), 
                         Inches(0.5), Inches(0.5),
                         font_size=Pt(16), color=self.palette["white"], 
                         bold=True, alignment=PP_ALIGN.CENTER)
            
            # 步骤标题
            step_title = step if isinstance(step, str) else str(step)
            self._add_text(slide, step_title[:20], 
                         x + Inches(0.1), start_y + Inches(0.9), 
                         step_width - Inches(0.2), Inches(1.5),
                         font_size=Pt(14), color=self.palette["text"],
                         alignment=PP_ALIGN.CENTER)
            
            # 箭头连接（除了最后一个）
            if idx < len(steps) - 1:
                arrow_x = x + step_width + Inches(0.05)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    arrow_x, start_y + Inches(1.1), Inches(0.2), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.palette["accent"]
                arrow.line.fill.background()
    
    def _render_list_content(self, slide, content: dict):
        """渲染列表类内容"""
        items = content.get("items", [])
        if not items:
            items = list(content.values())[:6] if isinstance(content, dict) else []
        
        # 计算布局
        cols = 3 if len(items) > 3 else len(items)
        rows = (len(items) + cols - 1) // cols
        
        card_width = Inches(3.8)
        card_height = Inches(2)
        start_x = Inches(0.8)
        start_y = Inches(1.3)
        
        for idx, item in enumerate(items):
            row = idx // cols
            col = idx % cols
            x = start_x + col * (card_width + Inches(0.3))
            y = start_y + row * (card_height + Inches(0.3))
            
            # 卡片
            self._add_card(slide, x, y, card_width, card_height, 
                          self.palette["white"])
            
            # 顶部色条
            self._add_gradient_bg(slide, x, y, card_width, Inches(0.12),
                                self.palette["primary"], self.palette["secondary"])
            
            # 内容
            item_text = item if isinstance(item, str) else str(item)
            self._add_text(slide, item_text[:50], 
                         x + Inches(0.2), y + Inches(0.3), 
                         card_width - Inches(0.4), Inches(1.5),
                         font_size=Pt(14), color=self.palette["text"])
    
    def _render_default_content(self, slide, content: dict):
        """默认渲染"""
        if isinstance(content, dict):
            items = list(content.values())
        elif isinstance(content, list):
            items = content
        else:
            items = [str(content)]
        
        y = Inches(1.5)
        for item in items[:8]:
            text = item if isinstance(item, str) else str(item)
            self._add_text(slide, f"• {text}", 
                         Inches(1), y, Inches(11), Inches(0.6),
                         font_size=Pt(16), color=self.palette["text"])
            y += Inches(0.7)
    
    def add_ending(self, text: str = "感谢聆听"):
        """添加结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self._add_gradient_bg(slide, 0, 0, 13.333, 7.5,
                              self.palette["primary"], self.palette["secondary"], 135)
        
        self._add_circle(slide, Inches(1), Inches(1), Inches(1.5), 
                        self.palette["accent"])
        self._add_circle(slide, Inches(10.5), Inches(5), Inches(2), 
                        self.palette["secondary"])
        
        self._add_text(slide, text, Inches(1), Inches(2.8), Inches(11), Inches(1),
                      font_size=Pt(40), color=self.palette["white"], bold=True,
                      alignment=PP_ALIGN.CENTER)
    
    # ==================== 基础组件 ====================
    
    def _add_gradient_bg(self, slide, left, top, width, height, color1, color2, angle=0):
        """添加渐变背景"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), 
            Inches(width), Inches(height)
        )
        shape.fill.background()
        fill = shape.fill
        fill.gradient()
        fill.gradient_angle = angle
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
        shape.line.fill.background()
    
    def _add_circle(self, slide, left, top, size, color):
        """添加圆形"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, size, size
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
    
    def _add_card(self, slide, left, top, width, height, color):
        """添加卡片"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = self.palette["light"]
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.05
    
    def _add_text(self, slide, text, left, top, width, height,
                 font_size=Pt(14), color=None, bold=False, alignment=PP_ALIGN.LEFT):
        """添加文本框"""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = alignment
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color or self.palette["text"]
    
    def save(self, path: str):
        """保存 PPT"""
        self.prs.save(path)
        print(f"✅ 智能 PPT 已生成: {path}")


# ==================== 演示函数 ====================
def demo():
    """演示 - 展示智能判断能力"""
    
    # 示例1：疗愈主题 - 自动匹配疗愈风配色
    generator1 = SmartPPTGerator("情感内耗疗愈沙龙")
    generator1.add_cover("情感内耗疗愈沙龙", "重新找回内心的平静")
    
    generator1.add_content_slide("什么是情感内耗", {
        "items": [
            "反复思考同一件事，无法停止",
            "过度在意他人看法和评价", 
            "情绪波动大，容易焦虑",
            "对自己要求过高，追求完美"
        ]
    })
    
    generator1.add_content_slide("疗愈四步法", {
        "steps": ["觉知", "接纳", "表达", "放下"]
    })
    
    generator1.add_content_slide("2024年服务数据", {
        "chart_data": {
            "categories": ["塔罗咨询", "占星服务", "SRT疗愈", "课程收入"],
            "series": {"营收占比": [35, 25, 20, 20]}
        }
    })
    
    generator1.add_ending("感谢聆听 · 重新拥抱自己")
    generator1.save("/Users/huanxi/Desktop/智能PPT_疗愈主题.pptx")
    
    # 示例2：商务主题 - 自动匹配商务风配色
    generator2 = SmartPPTGerator("财富能量沙龙")
    generator2.add_cover("财富能量公开课", "从能量层面提升财富频率")
    
    generator2.add_content_slide("今日大纲", {
        "items": [
            "财富能量的底层逻辑",
            "金钱限制性信念清理",
            "高频财富能量提升法",
            "现场财富疗愈体验"
        ]
    })
    
    generator2.add_content_slide("财富疗愈三阶段", {
        "steps": ["清理限制性信念", "提升能量频率", "显化财富目标"]
    })
    
    generator2.add_ending("开启你的财富之旅")
    generator2.save("/Users/huanxi/Desktop/智能PPT_商务主题.pptx")
    
    print("\n🎉 演示完成！生成了2个不同风格的智能 PPT")
    return True


if __name__ == "__main__":
    demo()
